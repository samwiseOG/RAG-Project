
import os
from pathlib import Path
from abc import ABC, abstractmethod

# Third-party libraries
import docx
import pdfplumber
import pptx



# 1. The Strategy Interface
class TextExtractor(ABC):
    """
    Declares an interface common to all supported extraction algorithms.
    The client code will use this interface to call the algorithm defined by a
    Concrete Strategy.
    """
    @abstractmethod
    def extract(self, file_path: str) -> str:
        """Extracts text content from a given file."""
        pass

# 2. Concrete Strategies
class DocxExtractor(TextExtractor):
    """Implements the text extraction algorithm for .docx files."""
    def extract(self, file_path: str) -> str:
        try:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            return f"Error processing DOCX file '{file_path}': {e}"

class PptxExtractor(TextExtractor):
    """Implements the text extraction algorithm for .pptx files."""
    def extract(self, file_path: str) -> str:
        try:
            presentation = pptx.Presentation(file_path)
            text_runs = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        text_runs.append(paragraph.text)
            return "\n".join(text_runs)
        except Exception as e:
            return f"Error processing PPTX file '{file_path}': {e}"

class PdfExtractor(TextExtractor):
    """Implements the text extraction algorithm for .pdf files."""
    def extract(self, file_path: str) -> str:
        try:
            with pdfplumber.open(file_path) as pdf:
                full_text = [p.extract_text() for p in pdf.pages if p.extract_text()]
            return "\n---\n".join(full_text)
        except Exception as e:
            return f"Error processing PDF file '{file_path}': {e}"

class TxtExtractor(TextExtractor):
    """Implements the text extraction algorithm for .txt files."""
    def extract(self, file_path: str) -> str:
        try:
            return Path(file_path).read_text(encoding='utf-8')
        except Exception as e:
            return f"Error processing TXT file '{file_path}': {e}"

class UnsupportedFileExtractor(TextExtractor):
    """A default strategy for unsupported file types."""
    def extract(self, file_path: str) -> str:
        file_extension = Path(file_path).suffix.lower()
        return f"Unsupported file type: '{file_extension}'"

# 3. The Simple Factory
def create_extractor(file_path: str) -> TextExtractor:
    """
    Factory function that selects and returns the appropriate extractor
    (Strategy) based on the file extension.
    """
    extension_map = {
        ".docx": DocxExtractor,
        ".pptx": PptxExtractor,
        ".pdf": PdfExtractor,
        ".txt": TxtExtractor,
    }
    file_extension = Path(file_path).suffix.lower()
    # Get the class from the map, or default to UnsupportedFileExtractor
    extractor_class = extension_map.get(file_extension, UnsupportedFileExtractor)
    return extractor_class()
